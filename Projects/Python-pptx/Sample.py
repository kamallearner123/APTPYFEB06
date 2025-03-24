from pptx import Presentation
from pptx.util import Inches

# Create a presentation
ppt = Presentation()

# Add Title Slide
slide_layout = ppt.slide_layouts[0]
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Django with Python"
subtitle.text = "An Interactive Guide"

# Add Content Slide
slide_layout = ppt.slide_layouts[1]
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "What is Django?"
content.text = "Django is a high-level Python web framework."

# Add Image Django.webp
slide_layout = ppt.slide_layouts[5]
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Django Logo"
img_path = "Django.png"
left = Inches(1)
top = Inches(1)
height = Inches(5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)


# Save the presentation
ppt.save("Django_Presentation.pptx")
